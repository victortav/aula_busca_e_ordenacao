"""
Algoritmos de Busca e Ordenacao - Slides v2
Melhorias: capa visual, logo IFPB, código colorido, slides de razões 1/4/5, slide mercado melhorado
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from PIL import Image, ImageDraw, ImageFont
import io, math, base64, textwrap

# ════════════════════════════════════════════════════════════════════════
# PALETTE  —  Ocean-Tech  (dark teal + cyan + midnight)
# ════════════════════════════════════════════════════════════════════════
C_BG_DARK  = RGBColor(0x06, 0x1A, 0x2E)   # almost-black navy
C_BG_MID   = RGBColor(0x06, 0x5A, 0x82)   # ocean blue (dark slide)
C_BG_LIGHT = RGBColor(0xF0, 0xF7, 0xFF)   # near-white (light slides)
C_TEAL     = RGBColor(0x00, 0xB4, 0xD8)   # vivid cyan accent
C_TEAL2    = RGBColor(0x1C, 0x72, 0x93)   # medium teal
C_MID      = RGBColor(0x03, 0x04, 0x5E)   # midnight indigo
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_GOLD     = RGBColor(0xFF, 0xC3, 0x00)   # amber highlight
C_GREEN    = RGBColor(0x06, 0xD6, 0xA0)   # emerald
C_RED      = RGBColor(0xEF, 0x47, 0x6F)   # coral red
C_ORANGE   = RGBColor(0xF7, 0x96, 0x1D)   # warm orange
C_GRAY     = RGBColor(0x64, 0x74, 0x8B)
C_DARK_TXT = RGBColor(0x1E, 0x29, 0x3B)   # body text on light bg

# Code colours
C_CODE_BG      = RGBColor(0x0D, 0x17, 0x26)
C_CODE_KEYWORD = RGBColor(0x56, 0x9C, 0xD6)  # blue
C_CODE_FUNC    = RGBColor(0xDC, 0xDC, 0xAA)  # gold-ish
C_CODE_STRING  = RGBColor(0xCE, 0x91, 0x78)  # salmon
C_CODE_COMMENT = RGBColor(0x6A, 0x99, 0x55)  # green
C_CODE_NUMBER  = RGBColor(0xB5, 0xCE, 0xA8)  # light green
C_CODE_DEFAULT = RGBColor(0xD4, 0xD4, 0xD4)  # light gray
C_CODE_HL      = RGBColor(0xA8, 0xFF, 0x78)  # bright green highlight

FONT_HEAD = "Trebuchet MS"
FONT_BODY = "Calibri"
FONT_CODE = "Consolas"

W_IN = 10.0
H_IN = 5.625
W = Inches(W_IN)
H = Inches(H_IN)

# ════════════════════════════════════════════════════════════════════════
# LOGO IFPB  —  drawn with Pillow, stored as PNG bytes
# ════════════════════════════════════════════════════════════════════════
def make_ifpb_logo(width=320, height=100):
    """Generate a clean IFPB text-badge logo."""
    img = Image.new("RGBA", (width, height), (0,0,0,0))
    d = ImageDraw.Draw(img)
    # background pill
    d.rounded_rectangle([0,0,width-1,height-1], radius=12,
                         fill=(0,180,216,255))  # teal
    # White bold text
    # Use default font since custom fonts may not be available
    try:
        font_big = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 44)
        font_sm  = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 18)
    except:
        font_big = ImageFont.load_default()
        font_sm  = font_big

    d.text((width//2, 38), "IFPB", font=font_big, fill=(255,255,255,255), anchor="mm")
    d.text((width//2, 76), "Instituto Federal da Paraíba", font=font_sm,
           fill=(220,240,255,220), anchor="mm")

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf

IFPB_LOGO = make_ifpb_logo()

# ════════════════════════════════════════════════════════════════════════
# COVER BACKGROUND  —  abstract circuit/grid pattern
# ════════════════════════════════════════════════════════════════════════
def make_cover_bg(w=1000, h=563):
    img = Image.new("RGB", (w, h), (6, 26, 46))
    d = ImageDraw.Draw(img)

    # Grid dots
    for gx in range(0, w, 40):
        for gy in range(0, h, 40):
            d.ellipse([gx-1, gy-1, gx+1, gy+1], fill=(0,100,140,255))

    # Diagonal accent lines
    for i in range(0, w+h, 80):
        x0 = i; y0 = 0; x1 = i-h; y1 = h
        d.line([x0,y0,x1,y1], fill=(0,90,130), width=1)

    # Glowing circles (large)
    for cx, cy, r, col in [
        (720, 280, 200, (0,180,216,30)),
        (820, 100, 120, (0,100,180,25)),
        (600, 420, 90,  (0,150,180,20)),
    ]:
        for dr in range(r, 0, -20):
            alpha = int(col[3] * (1 - dr/r))
            d.ellipse([cx-dr, cy-dr, cx+dr, cy+dr],
                      outline=(col[0],col[1],col[2],max(5,alpha)), width=1)

    # Bright horizontal accent band
    for y in range(260, 310):
        alpha = int(60 * math.sin(math.pi * (y-260)/50))
        d.line([0,y,w,y], fill=(0,180,216, alpha))

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf

COVER_BG = make_cover_bg()

# ════════════════════════════════════════════════════════════════════════
# HELPERS
# ════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

def slide(dark=True):
    s = prs.slides.add_slide(BLANK)
    if dark:
        add_rect(s, 0, 0, W_IN, H_IN, fill=C_BG_DARK)
    else:
        add_rect(s, 0, 0, W_IN, H_IN, fill=C_BG_LIGHT)
    return s

def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_w=Pt(0), radius=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_w or Pt(1)
    else:
        shape.line.fill.background()
    return shape

def txt(s, text, x, y, w, h, size=16, bold=False, color=C_WHITE,
        align=PP_ALIGN.LEFT, font=FONT_BODY, italic=False):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb

def txt_ml(s, lines, x, y, w, h, size=14, color=C_WHITE, font=FONT_BODY, spacing=Pt(4)):
    """Multi-line: lines = str | (text, bold, color)"""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            t, b, c = item, False, color
        else:
            t, b, c = item[0], item[1], (item[2] if len(item)>2 else color)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.space_after = spacing
        r = p.add_run(); r.text = t
        r.font.name = font; r.font.size = Pt(size)
        r.font.bold = b; r.font.color.rgb = c
        first = False
    return tb

def add_logo(s, x=8.7, y=0.08, w=1.1, h=0.35):
    """Place IFPB logo pill on slide (top-right area)."""
    IFPB_LOGO.seek(0)
    s.shapes.add_picture(IFPB_LOGO, Inches(x), Inches(y), Inches(w), Inches(h))

def section_tag(s, label, x=0.35, y=0.08):
    """Small section label top-left."""
    add_rect(s, x, y, 2.6, 0.28, fill=C_TEAL2)
    txt(s, label, x+0.08, y+0.01, 2.5, 0.26, size=11, bold=True,
        color=C_WHITE, font=FONT_HEAD)

# ────────────────────────────────────────────────────────────────────────
# CODE BLOCK with syntax colouring  (python-pptx rich text)
# ────────────────────────────────────────────────────────────────────────
KEYWORDS = {"def","for","while","if","elif","else","return","import","from",
            "in","not","and","or","True","False","None","class","with","as",
            "try","except","pass","lambda","yield","raise"}
BUILTINS = {"len","range","print","enumerate","sorted","min","max","list",
            "dict","set","int","str","float","bool","abs","sum","type"}

def _tokenize(line):
    """Very light tokenizer → list of (text, style)
       style: 'keyword','builtin','string','comment','number','default'
    """
    if line.strip().startswith("#"):
        return [(line, "comment")]
    tokens = []
    i = 0
    while i < len(line):
        c = line[i]
        # comment mid-line
        if c == "#":
            tokens.append((line[i:], "comment"))
            break
        # string
        if c in ('"', "'"):
            q = c
            j = i+1
            while j < len(line) and line[j] != q:
                j += 1
            tokens.append((line[i:j+1], "string"))
            i = j+1
            continue
        # word
        if c.isalpha() or c == "_":
            j = i
            while j < len(line) and (line[j].isalnum() or line[j] == "_"):
                j += 1
            word = line[i:j]
            if word in KEYWORDS:
                tokens.append((word, "keyword"))
            elif word in BUILTINS:
                tokens.append((word, "builtin"))
            else:
                tokens.append((word, "default"))
            i = j
            continue
        # number
        if c.isdigit():
            j = i
            while j < len(line) and (line[j].isdigit() or line[j] == "."):
                j += 1
            tokens.append((line[i:j], "number"))
            i = j
            continue
        tokens.append((c, "default"))
        i += 1
    return tokens

TOKEN_COLOR = {
    "keyword":  C_CODE_KEYWORD,
    "builtin":  C_CODE_FUNC,
    "string":   C_CODE_STRING,
    "comment":  C_CODE_COMMENT,
    "number":   C_CODE_NUMBER,
    "default":  C_CODE_DEFAULT,
}

def code_block(s, lines, x, y, w, h, font_size=11.5, highlight_lines=None):
    """Draw a syntax-coloured code block."""
    highlight_lines = highlight_lines or []
    add_rect(s, x, y, w, h, fill=C_CODE_BG)
    # Thin teal top-border accent
    add_rect(s, x, y, w, 0.04, fill=C_TEAL)

    tb = s.shapes.add_textbox(Inches(x+0.15), Inches(y+0.10),
                               Inches(w-0.3), Inches(h-0.18))
    tb.word_wrap = False
    tf = tb.text_frame; tf.word_wrap = False

    first = True
    for li, line in enumerate(lines):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.space_after = Pt(1)

        if not line.strip():
            r = p.add_run(); r.text = " "
            r.font.name = FONT_CODE; r.font.size = Pt(font_size)
            r.font.color.rgb = C_CODE_DEFAULT
            first = False
            continue

        tokens = _tokenize(line)
        for tok_text, tok_style in tokens:
            r = p.add_run()
            r.text = tok_text
            r.font.name = FONT_CODE
            r.font.size = Pt(font_size)
            r.font.color.rgb = TOKEN_COLOR.get(tok_style, C_CODE_DEFAULT)
            # Highlight entire line
            if li in highlight_lines:
                r.font.color.rgb = C_CODE_HL

        first = False

# ────────────────────────────────────────────────────────────────────────
# Info / callout card
# ────────────────────────────────────────────────────────────────────────
def callout(s, title, lines, x, y, w, h, accent=C_TEAL, bg=None, title_size=13, body_size=12):
    bg = bg or C_BG_MID
    add_rect(s, x, y, w, h, fill=bg)
    add_rect(s, x, y, 0.05, h, fill=accent)
    txt(s, title, x+0.15, y+0.08, w-0.2, 0.32,
        size=title_size, bold=True, color=accent, font=FONT_HEAD)
    txt_ml(s, lines, x+0.15, y+0.42, w-0.25, h-0.5,
           size=body_size, color=C_WHITE)

def complexity_badge(s, notation, color, x, y):
    add_rect(s, x, y, 1.35, 0.38, fill=color)
    txt(s, notation, x, y, 1.35, 0.38, size=14, bold=True,
        color=C_BG_DARK, align=PP_ALIGN.CENTER, font=FONT_CODE)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 1  —  COVER  (redesigned with image + vivid elements)
# ════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
# Full-bleed circuit background
COVER_BG.seek(0)
s1.shapes.add_picture(COVER_BG, Inches(0), Inches(0), Inches(W_IN), Inches(H_IN))

# Dark left panel
add_rect(s1, 0, 0, 5.8, H_IN, fill=C_BG_DARK)
# Vivid teal vertical accent
add_rect(s1, 5.65, 0, 0.12, H_IN, fill=C_TEAL)

# IFPB logo top-right
IFPB_LOGO.seek(0)
s1.shapes.add_picture(IFPB_LOGO, Inches(6.5), Inches(0.2), Inches(1.6), Inches(0.5))

# Tagline chip
add_rect(s1, 0.5, 0.55, 2.4, 0.32, fill=C_TEAL)
txt(s1, "CURSO TÉCNICO DE PROGRAMAÇÃO", 0.5, 0.55, 2.4, 0.32,
    size=8.5, bold=True, color=C_BG_DARK, align=PP_ALIGN.CENTER, font=FONT_HEAD)

# Big title
txt(s1, "ALGORITMOS", 0.5, 1.05, 5.0, 0.82,
    size=44, bold=True, color=C_WHITE, font=FONT_HEAD)
txt(s1, "DE BUSCA &", 0.5, 1.82, 5.0, 0.75,
    size=44, bold=True, color=C_WHITE, font=FONT_HEAD)
txt(s1, "ORDENAÇÃO", 0.5, 2.52, 5.0, 0.78,
    size=44, bold=True, color=C_TEAL, font=FONT_HEAD)

# Subtitle
txt(s1, "Do conceito ao código — o que você vai usar na prática",
    0.5, 3.42, 5.1, 0.5, size=14, italic=True,
    color=RGBColor(0xA0,0xCC,0xE8), font=FONT_BODY)

# Bottom info bar
add_rect(s1, 0, 5.05, 5.8, 0.575, fill=C_MID)
txt(s1, "Aula de 40 min  ·  Python 3  ·  IFPB — 2026",
    0.5, 5.1, 5.1, 0.45, size=12, color=RGBColor(0x90,0xB8,0xD0), font=FONT_BODY)

# Right side: 3 floating stat chips
chips = [
    ("4", "algoritmos"),
    ("O(log n)", "eficiência"),
    ("100%", "prático"),
]
for i, (big, small) in enumerate(chips):
    cy = 1.2 + i * 1.35
    add_rect(s1, 6.3, cy, 3.3, 1.05,
             fill=RGBColor(0x06, 0x5A, 0x82))
    add_rect(s1, 6.3, cy, 3.3, 0.07, fill=C_TEAL)
    txt(s1, big,   6.35, cy+0.1,  3.1, 0.55,
        size=32, bold=True, color=C_WHITE, font=FONT_HEAD, align=PP_ALIGN.CENTER)
    txt(s1, small, 6.35, cy+0.62, 3.1, 0.35,
        size=13, color=C_TEAL, font=FONT_BODY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 2  —  MOTIVAÇÃO
# ════════════════════════════════════════════════════════════════════════
s2 = slide(dark=False)
add_logo(s2)
section_tag(s2, "INTRODUÇÃO  ·  Por que isso importa?")

txt(s2, "Imagine este problema...", 0.35, 0.45, 9, 0.55,
    size=22, bold=True, color=C_BG_DARK, font=FONT_HEAD)

# Left card
add_rect(s2, 0.35, 1.1, 4.3, 1.95, fill=C_WHITE,
         line_color=C_TEAL2, line_w=Pt(1.5))
add_rect(s2, 0.35, 1.1, 4.3, 0.38, fill=C_TEAL2)
txt(s2, "📋  Lista de chamada", 0.5, 1.13, 4.0, 0.32,
    size=13, bold=True, color=C_WHITE, font=FONT_HEAD)
txt_ml(s2, [
    "500 alunos em ordem aleatória.",
    "Professor precisa achar \"Maria Silva\".",
    "⏱  Quanto tempo leva?",
], 0.5, 1.55, 4.0, 1.4, size=13, color=C_DARK_TXT)

# Right card
add_rect(s2, 5.2, 1.1, 4.45, 1.95, fill=C_WHITE,
         line_color=C_GREEN, line_w=Pt(1.5))
add_rect(s2, 5.2, 1.1, 4.45, 0.38, fill=C_GREEN)
txt(s2, "📚  Lista em ordem alfabética", 5.35, 1.13, 4.1, 0.32,
    size=13, bold=True, color=C_WHITE, font=FONT_HEAD)
txt_ml(s2, [
    "Mesmos 500 alunos — ordenados.",
    'Vai direto para a letra "M".',
    "⚡  Até 100× mais rápido!",
], 5.35, 1.55, 4.1, 1.4, size=13, color=C_DARK_TXT)

txt(s2, "👉  Ordenar e buscar dados é uma das tarefas mais comuns em qualquer sistema real.",
    0.35, 3.18, 9.3, 0.42, size=14, bold=True, color=C_BG_MID, font=FONT_BODY)

txt_ml(s2, [
    ("Exemplos do dia a dia do dev:", True, C_BG_DARK),
    "  Ordenar resultados por relevância  ·  Pesquisar usuário no banco  ·  Autocompletar  ·  Ranking por preço",
], 0.35, 3.68, 9.3, 1.0, size=13, color=C_DARK_TXT)
add_logo(s2)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 3  —  ROTEIRO
# ════════════════════════════════════════════════════════════════════════
s3 = slide(dark=True)
add_logo(s3)
section_tag(s3, "ROTEIRO  ·  40 minutos")

txt(s3, "O que vamos aprender hoje", 0.35, 0.45, 9, 0.55,
    size=24, bold=True, color=C_WHITE, font=FONT_HEAD)

items = [
    ("01", "Por que busca e ordenação são cruciais", "5 min", C_TEAL),
    ("02", "Busca Linear vs Busca Binária — com analogia", "8 min", C_GREEN),
    ("03", "Bubble Sort vs Merge Sort — do ingênuo ao eficiente", "10 min", C_GOLD),
    ("04", "Complexidade de Tempo — leitura rápida do Big O", "7 min", C_ORANGE),
    ("05", "Código Python na prática + quadro branco", "7 min", C_TEAL),
    ("06", "O que o mercado usa — dica de carreira", "3 min", C_GREEN),
]
for i, (num, desc, dur, col) in enumerate(items):
    yy = 1.2 + i * 0.68
    add_rect(s3, 0.35, yy, 0.52, 0.48, fill=col)
    txt(s3, num, 0.35, yy, 0.52, 0.48, size=16, bold=True,
        color=C_BG_DARK, align=PP_ALIGN.CENTER, font=FONT_HEAD)
    txt(s3, desc, 1.0, yy+0.06, 7.0, 0.36, size=14, color=C_WHITE, font=FONT_BODY)
    add_rect(s3, 8.2, yy, 1.35, 0.48, fill=RGBColor(0x06,0x3A,0x56))
    txt(s3, dur, 8.2, yy, 1.35, 0.48, size=13, color=col,
        align=PP_ALIGN.CENTER, font=FONT_BODY, bold=True)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 4  —  BUSCA LINEAR
# ════════════════════════════════════════════════════════════════════════
s4 = slide(dark=False)
add_logo(s4)
section_tag(s4, "BUSCA  ·  Linear Search")

txt(s4, "Busca Linear — o jeito mais simples (e mais lento)", 0.35, 0.45, 9, 0.5,
    size=20, bold=True, color=C_BG_DARK, font=FONT_HEAD)
txt(s4, "🎲  Analogia: procurar um número numa fila de caixas fechadas",
    0.35, 1.0, 9, 0.36, size=13, italic=True, color=C_TEAL2)

boxes = [7, 2, 14, 5, 9, 11, 3]
target_idx = 4
bx = 0.35
for i, val in enumerate(boxes):
    bg = C_GOLD if i == target_idx else (C_TEAL2 if i < target_idx else C_GRAY)
    add_rect(s4, bx + i*1.3, 1.42, 1.1, 0.78, fill=bg)
    txt(s4, str(val), bx + i*1.3, 1.42, 1.1, 0.78,
        size=22, bold=True,
        color=C_BG_DARK if i == target_idx else C_WHITE,
        align=PP_ALIGN.CENTER, font=FONT_HEAD)
    icon = "✅" if i == target_idx else ("👁" if i < target_idx else "")
    if icon:
        txt(s4, icon, bx + i*1.3, 2.24, 1.1, 0.28,
            size=13, align=PP_ALIGN.CENTER, color=C_DARK_TXT)

txt(s4, "Procurando o 9: verificou 7→2→14→5 antes de achar — percorreu 5 elementos",
    0.35, 2.58, 9.3, 0.32, size=12, color=C_GRAY, italic=True)

code_block(s4, [
    "def busca_linear(lista, alvo):",
    "    for i, valor in enumerate(lista):",
    "        if valor == alvo:",
    "            return i       # índice encontrado",
    "    return -1              # não encontrado",
], 0.35, 3.0, 4.8, 1.9, highlight_lines=[2,3])

callout(s4, "Complexidade", [
    ("Melhor caso:", True, C_GREEN),
    ("  O(1) — achou na 1ª posição", False, C_WHITE),
    ("Pior / médio caso:", True, C_ORANGE),
    ("  O(n) — percorre toda a lista", False, C_WHITE),
    ("Quando usar:", True, C_TEAL),
    ("  Lista pequena ou não ordenada", False, C_WHITE),
], 5.4, 3.0, 4.25, 1.9, accent=C_TEAL, bg=C_BG_MID)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 5  —  BUSCA BINÁRIA
# ════════════════════════════════════════════════════════════════════════
s5 = slide(dark=False)
add_logo(s5)
section_tag(s5, "BUSCA  ·  Binary Search")

txt(s5, "Busca Binária — o salto quântico de eficiência", 0.35, 0.45, 9, 0.5,
    size=20, bold=True, color=C_BG_DARK, font=FONT_HEAD)
txt(s5, "📖  Analogia: como você procura uma palavra no dicionário?",
    0.35, 1.0, 9, 0.36, size=13, italic=True, color=C_TEAL2)

sorted_vals = [2, 5, 7, 9, 11, 14, 18]
bx = 0.35
for i, val in enumerate(sorted_vals):
    bg = C_GOLD if i == 3 else RGBColor(0xBB,0xCC,0xDD)
    add_rect(s5, bx + i*1.3, 1.42, 1.1, 0.68, fill=bg)
    txt(s5, str(val), bx + i*1.3, 1.42, 1.1, 0.68,
        size=20, bold=True, color=C_BG_DARK, align=PP_ALIGN.CENTER, font=FONT_HEAD)

txt(s5, "Meio ↑", bx+3*1.3, 2.14, 1.1, 0.28,
    size=11, color=C_GOLD, bold=True, align=PP_ALIGN.CENTER)
txt(s5, "Alvo = 14 > 9  →  descarta metade esquerda  →  próxima busca: {11, 14, 18}",
    0.35, 2.48, 9.3, 0.32, size=12, italic=True, color=C_GRAY)

code_block(s5, [
    "def busca_binaria(lista, alvo):",
    "    esq, dir = 0, len(lista) - 1",
    "    while esq <= dir:",
    "        meio = (esq + dir) // 2",
    "        if lista[meio] == alvo:",
    "            return meio         # ✅ encontrado!",
    "        elif lista[meio] < alvo:",
    "            esq = meio + 1      # descarta esquerda",
    "        else:",
    "            dir = meio - 1      # descarta direita",
    "    return -1",
], 0.35, 2.88, 5.0, 2.55)

callout(s5, "Por que é tão rápido?", [
    ("O(log n)", True, C_GREEN),
    ("", False, C_WHITE),
    ("  1.000 itens → máx. 10 comp.", False, C_WHITE),
    ("  1.000.000 → máx. 20 comp.!", False, C_WHITE),
    ("", False, C_WHITE),
    ("Pré-requisito:", True, C_ORANGE),
    ("  Lista ORDENADA", False, C_WHITE),
    ("", False, C_WHITE),
    ("Na prática:", True, C_TEAL),
    ("  Python: bisect.bisect_left()", False, C_CODE_FUNC),
    ("  Java: Arrays.binarySearch()", False, C_CODE_FUNC),
], 5.55, 2.88, 4.1, 2.55, accent=C_GREEN, bg=C_BG_MID)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 6  —  BUBBLE SORT
# ════════════════════════════════════════════════════════════════════════
s6 = slide(dark=False)
add_logo(s6)
section_tag(s6, "ORDENAÇÃO  ·  Bubble Sort — o ingênuo (mas didático)")

txt(s6, "Bubble Sort — entendendo o problema antes da solução", 0.35, 0.45, 9, 0.5,
    size=19, bold=True, color=C_BG_DARK, font=FONT_HEAD)
txt(s6, "🫧  Analogia: bolhas subindo — números maiores 'flutuam' para o fim",
    0.35, 0.98, 9, 0.34, size=13, italic=True, color=C_TEAL2)

steps = [
    ([5, 3, 8, 1], "Original"),
    ([3, 5, 1, 8], "Passo 1 (8↑)"),
    ([3, 1, 5, 8], "Passo 2 (5↑)"),
    ([1, 3, 5, 8], "Ordenado! ✅"),
]
scols = [
    [C_GRAY, C_GRAY, C_GOLD, C_GRAY],
    [C_GRAY, C_GRAY, C_GOLD, C_GREEN],
    [C_GRAY, C_GRAY, C_GOLD, C_GREEN],
    [C_GREEN]*4,
]
for si, (vals, label) in enumerate(steps):
    bxs = 0.3 + si * 2.38
    txt(s6, label, bxs, 1.38, 2.1, 0.28, size=11,
        color=C_GRAY, align=PP_ALIGN.CENTER)
    for j, v in enumerate(vals):
        add_rect(s6, bxs + j*0.52, 1.68, 0.46, 0.58, fill=scols[si][j])
        txt(s6, str(v), bxs + j*0.52, 1.68, 0.46, 0.58,
            size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

code_block(s6, [
    "def bubble_sort(lista):",
    "    n = len(lista)",
    "    for i in range(n):",
    "        for j in range(0, n - i - 1):",
    "            if lista[j] > lista[j+1]:",
    "                # troca os elementos",
    "                lista[j], lista[j+1] = lista[j+1], lista[j]",
    "    return lista",
], 0.35, 2.42, 5.0, 2.08)

callout(s6, "⚠️  Por que NÃO usar em produção", [
    ("O(n²)", True, C_RED),
    ("", False, C_WHITE),
    ("  10.000 itens → 100M compar. 😱", False, C_WHITE),
    ("", False, C_WHITE),
    ("Use para:", True, C_GOLD),
    ("  Aprendizado e entrevistas.", False, C_WHITE),
    ("  NUNCA em dados reais grandes.", False, C_WHITE),
], 5.55, 2.42, 4.1, 2.08, accent=C_RED, bg=RGBColor(0x2A,0x10,0x10))

txt(s6, "💡  Entender o Bubble Sort revela POR QUÊ algoritmos melhores existem.",
    0.35, 4.6, 9.3, 0.36, size=12.5, bold=True, color=C_TEAL2, italic=True)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 7  —  MERGE SORT
# ════════════════════════════════════════════════════════════════════════
s7 = slide(dark=False)
add_logo(s7)
section_tag(s7, "ORDENAÇÃO  ·  Merge Sort — dividir para conquistar")

txt(s7, "Merge Sort — a estratégia inteligente", 0.35, 0.45, 9, 0.5,
    size=20, bold=True, color=C_BG_DARK, font=FONT_HEAD)
txt(s7, "🃏  Analogia: embaralhou dois baralhos já ordenados — é fácil intercalá-los!",
    0.35, 0.98, 9, 0.34, size=13, italic=True, color=C_TEAL2)

def ms_box(sl, x, y, vals, col=C_TEAL2):
    w = len(vals) * 0.42 + 0.1
    add_rect(sl, x, y, w, 0.42, fill=col)
    for k, v in enumerate(vals):
        txt(sl, str(v), x+k*0.42, y, 0.42, 0.42,
            size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

ms_box(s7, 1.85, 1.4, [8,3,1,5,7,2])
ms_box(s7, 0.4, 2.1, [8,3,1], col=C_MID)
ms_box(s7, 4.8, 2.1, [5,7,2], col=C_MID)
ms_box(s7, 0.4, 2.82, [8], col=C_GREEN)
ms_box(s7, 1.28, 2.82, [3,1], col=C_GREEN)
ms_box(s7, 4.8, 2.82, [5], col=C_GREEN)
ms_box(s7, 5.68, 2.82, [7,2], col=C_GREEN)
ms_box(s7, 1.85, 3.58, [1,2,3,5,7,8], col=C_TEAL)

txt(s7, "⬇ divide", 3.52, 2.2, 1.0, 0.28, size=11, color=C_GRAY, align=PP_ALIGN.CENTER)
txt(s7, "⬇ divide", 3.52, 2.92, 1.0, 0.28, size=11, color=C_GRAY, align=PP_ALIGN.CENTER)
txt(s7, "⬆ merge", 3.52, 3.38, 1.0, 0.28, size=11, color=C_GREEN, align=PP_ALIGN.CENTER)

code_block(s7, [
    "def merge_sort(lista):",
    "    if len(lista) <= 1:",
    "        return lista",
    "    meio = len(lista) // 2",
    "    esq = merge_sort(lista[:meio])",
    "    dir = merge_sort(lista[meio:])",
    "    return merge(esq, dir)  # intercala",
], 6.1, 2.1, 3.55, 2.3)

add_rect(s7, 6.1, 4.5, 3.55, 0.7, fill=RGBColor(0x06,0x3A,0x20))
txt(s7, "✅  O(n log n)  — Python sorted() usa Timsort, variante do Merge Sort",
    6.25, 4.6, 3.3, 0.5, size=12, bold=True, color=C_GREEN, font=FONT_BODY)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 8  —  BIG O
# ════════════════════════════════════════════════════════════════════════
s8 = slide(dark=True)
add_logo(s8)
section_tag(s8, "COMPLEXIDADE  ·  Leitura rápida do Big O")

txt(s8, "Como comparar algoritmos sem nem rodar o código", 0.35, 0.45, 9, 0.5,
    size=20, bold=True, color=C_WHITE, font=FONT_HEAD)

headers = ["Notação", "Nome", "Exemplo", "n = 1.000"]
col_w   = [1.45, 1.55, 3.55, 2.2]
col_x   = [0.3, 1.8, 3.4, 7.0]
hdr_bg  = RGBColor(0x03,0x04,0x5E)

bdr = { style: BorderStyle.NONE for style in [] }  # no border placeholder

rows = [
    ["O(1)",       "Constante",   "Acesso em array por índice",   "1 op",            C_GREEN],
    ["O(log n)",   "Logarítmica", "Busca Binária",                "~10 ops",         RGBColor(0x06,0xA0,0x70)],
    ["O(n)",       "Linear",      "Busca Linear",                 "1.000 ops",       C_TEAL2],
    ["O(n log n)", "Log-linear",  "Merge Sort / Timsort",         "~10.000 ops",     C_MID],
    ["O(n²)",      "Quadrática",  "Bubble Sort",                  "1.000.000 ops 😱",C_RED],
]

for ci, (hd, cw, cx) in enumerate(zip(headers, col_w, col_x)):
    add_rect(s8, cx, 1.08, cw, 0.38, fill=hdr_bg)
    txt(s8, hd, cx+0.06, 1.08, cw-0.1, 0.38, size=13, bold=True,
        color=C_TEAL, align=PP_ALIGN.CENTER, font=FONT_HEAD)

rh = 0.52
for ri, (row_vals) in enumerate(rows):
    vals = row_vals[:-1]; rc = row_vals[-1]
    yy = 1.5 + ri * rh
    for ci, (cell, cw, cx) in enumerate(zip(vals, col_w, col_x)):
        bg = rc if ci == 0 else RGBColor(0x0F,0x2A,0x40)
        add_rect(s8, cx, yy, cw, rh-0.04, fill=bg)
        fc = C_BG_DARK if ci == 0 else (C_GOLD if "😱" in cell else RGBColor(0xCC,0xDD,0xFF))
        alg = PP_ALIGN.CENTER if ci != 2 else PP_ALIGN.LEFT
        txt(s8, cell, cx+0.06, yy+0.07, cw-0.12, rh-0.14,
            size=12.5, color=fc, bold=(ci==0),
            font=FONT_CODE if ci==0 else FONT_BODY, align=alg)

txt(s8, "💡  Regra prática: prefira O(log n) ou O(n log n) para listas grandes.",
    0.3, 5.15, 9.3, 0.32, size=13, bold=True, color=C_TEAL, font=FONT_BODY)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 9  —  EXERCÍCIO / QUADRO BRANCO
# ════════════════════════════════════════════════════════════════════════
s9 = slide(dark=True)
add_logo(s9)
section_tag(s9, "PRÁTICA  ·  Vamos ao quadro branco!")

txt(s9, "🖊️  Exercício em duplas — 5 minutos", 0.4, 0.45, 9, 0.5,
    size=22, bold=True, color=C_GOLD, font=FONT_HEAD)

add_rect(s9, 0.4, 1.08, 9.2, 1.38, fill=RGBColor(0x0F,0x2A,0x40))
txt(s9, "Dada a lista:   [42, 7, 15, 3, 28, 11, 19]",
    0.6, 1.16, 8.8, 0.42, size=18, bold=True, color=C_GOLD, font=FONT_CODE)
txt_ml(s9, [
    "a)  Execute o Bubble Sort — quais trocas ocorrem na 1ª passagem?",
    "b)  Use Busca Binária para encontrar o 19 (após ordenar). Quantas comparações?",
], 0.6, 1.6, 8.8, 0.76, size=14, color=RGBColor(0xCC,0xDD,0xFF))

add_rect(s9, 0.4, 2.55, 9.2, 0.04, fill=C_TEAL)

txt(s9, "📋  Resolução para o professor (quadro):", 0.5, 2.68, 9, 0.35,
    size=13, bold=True, color=C_TEAL, font=FONT_HEAD)
txt_ml(s9, [
    "a)  42>7 TROCA  →  42>15 TROCA  →  42>3 TROCA  →  42>28 TROCA  →  42>11 TROCA  →  42>19 TROCA  →  [ 7,15,3,28,11,19,42 ]",
    "b)  [ 3,7,11,15,19,28,42 ]  meio=15 → 19>15 → dir {19,28,42}  meio=28 → 19<28 → esq {19}  ✅  (3 comparações)",
], 0.5, 3.08, 9.2, 0.9, size=12.5, color=RGBColor(0xB0,0xD0,0xFF))

add_rect(s9, 0.4, 4.1, 9.2, 0.82, fill=RGBColor(0x06,0x3A,0x28))
txt(s9, "🔑  Insight:",
    0.6, 4.15, 1.2, 0.32, size=13, bold=True, color=C_GREEN, font=FONT_HEAD)
txt(s9, "Ordenar uma vez + Busca Binária = muito mais rápido que buscar repetidamente sem ordem.",
    1.7, 4.2, 7.7, 0.62, size=13, color=C_WHITE)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 10  —  MERCADO (melhorado)
# ════════════════════════════════════════════════════════════════════════
s10 = slide(dark=False)
add_logo(s10)
section_tag(s10, "MERCADO  ·  O que você vai usar como desenvolvedor")

txt(s10, "Ninguém implementa do zero — mas entender é essencial", 0.35, 0.45, 9.3, 0.5,
    size=19, bold=True, color=C_BG_DARK, font=FONT_HEAD)

cards = [
    ("Python", "sorted() / list.sort()", C_TEAL2,
     "Algoritmo: Timsort (Merge Sort + Insertion Sort)",
     "In-place: list.sort()  |  Nova lista: sorted()",
     "Estável ✅  |  O(n log n) garantido"),
    ("Java", "Arrays.sort() / Collections.sort()", C_MID,
     "Primitivos: Dual-Pivot Quicksort (O(n log n) médio)",
     "Objetos: Timsort (estável, O(n log n))",
     "In-place para arrays  |  Nova lista para Collections"),
    ("SQL", "ORDER BY / índices B-Tree", C_GREEN,
     "Algoritmo: B-Tree para buscas em índice",
     "Busca Binária na árvore — O(log n) por query",
     "Ordenação: depende do engine (QuickSort/MergeSort)"),
    ("Entrevistas", "LeetCode / HackerRank", C_ORANGE,
     "Você SERÁ perguntado sobre Big O",
     "Conhecer Merge Sort / Binary Search = diferencial",
     "Sênior vs Júnior começa aqui"),
]

for i, (lang, api, col, l1, l2, l3) in enumerate(cards):
    xi = 0.3 + (i % 2) * 4.88
    yi = 1.08 + (i // 2) * 2.15

    # Header
    add_rect(s10, xi, yi, 4.5, 0.46, fill=col)
    txt(s10, lang, xi+0.12, yi+0.04, 2.1, 0.38,
        size=15, bold=True, color=C_WHITE, font=FONT_HEAD)
    txt(s10, api, xi+2.35, yi+0.07, 2.1, 0.32,
        size=10.5, color=C_WHITE, font=FONT_CODE)

    # Body
    add_rect(s10, xi, yi+0.46, 4.5, 1.6, fill=C_WHITE,
             line_color=col, line_w=Pt(1))
    # Accent left border
    add_rect(s10, xi, yi+0.46, 0.05, 1.6, fill=col)

    txt_ml(s10, [
        (l1, True, C_BG_DARK),
        (l2, False, C_DARK_TXT),
        (l3, False, C_TEAL2),
    ], xi+0.15, yi+0.54, 4.25, 1.44, size=12, color=C_DARK_TXT)

txt(s10, "✅  Entender a complexidade dos algoritmos da stdlib torna você um dev muito mais sólido.",
    0.3, 5.2, 9.4, 0.28, size=12.5, bold=True, color=C_TEAL2, italic=True)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 11  —  RAZÃO 1: Escolher o algoritmo certo
# ════════════════════════════════════════════════════════════════════════
s11 = slide(dark=True)
add_logo(s11)
section_tag(s11, "POR QUE APRENDER  ·  Razão 1 de 3")

txt(s11, "Você escolhe a ferramenta certa para o problema certo", 0.35, 0.45, 9.3, 0.5,
    size=20, bold=True, color=C_WHITE, font=FONT_HEAD)
txt(s11, "A stdlib tem várias opções — mas cada uma tem um custo diferente.",
    0.35, 0.98, 9.3, 0.32, size=13, italic=True,
    color=RGBColor(0xA0,0xCC,0xE8))

# Left: bad vs good (k-smallest)
add_rect(s11, 0.35, 1.38, 4.6, 0.28, fill=RGBColor(0x6B,0x10,0x10))
txt(s11, "❌  O(n log n)  — ordena TUDO", 0.45, 1.38, 4.4, 0.28,
    size=11, bold=True, color=C_RED, font=FONT_HEAD)

code_block(s11, [
    "# Quer apenas os 5 menores de 10M números",
    "menores = sorted(dados)[:5]  # ordena tudo!",
], 0.35, 1.68, 4.6, 0.88)

add_rect(s11, 0.35, 2.6, 4.6, 0.28, fill=RGBColor(0x06,0x3A,0x20))
txt(s11, "✅  O(n log k)  — muito mais rápido", 0.45, 2.6, 4.4, 0.28,
    size=11, bold=True, color=C_GREEN, font=FONT_HEAD)

code_block(s11, [
    "import heapq",
    "# Heap busca os k menores sem ordenar tudo",
    "menores = heapq.nsmallest(5, dados)",
], 0.35, 2.9, 4.6, 1.0)

# Right: linear vs binary search
add_rect(s11, 5.1, 1.38, 4.55, 0.28, fill=RGBColor(0x6B,0x10,0x10))
txt(s11, "❌  O(n)  — busca linear em lista ordenada", 5.2, 1.38, 4.3, 0.28,
    size=11, bold=True, color=C_RED, font=FONT_HEAD)

code_block(s11, [
    "# Python não sabe que está ordenada!",
    "if 22 in lista_ordenada:  # O(n)",
    "    print('achou')",
], 5.1, 1.68, 4.55, 0.88)

add_rect(s11, 5.1, 2.6, 4.55, 0.28, fill=RGBColor(0x06,0x3A,0x20))
txt(s11, "✅  O(log n)  — usa bisect", 5.2, 2.6, 4.3, 0.28,
    size=11, bold=True, color=C_GREEN, font=FONT_HEAD)

code_block(s11, [
    "import bisect",
    "pos = bisect.bisect_left(lista_ordenada, 22)",
    "achou = (lista_ordenada[pos] == 22)",
    "# Você diz que está ordenada → O(log n)",
], 5.1, 2.9, 4.55, 1.0)

add_rect(s11, 0.35, 4.05, 9.3, 0.72, fill=RGBColor(0x06,0x3A,0x56))
txt(s11, "💡  A lib não decide por você. Só você sabe que a lista está ordenada — e pode usar Busca Binária.",
    0.5, 4.12, 9.0, 0.58, size=13, bold=True, color=C_TEAL, font=FONT_BODY)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 12  —  RAZÃO 4: Big O é a língua franca
# ════════════════════════════════════════════════════════════════════════
s12 = slide(dark=True)
add_logo(s12)
section_tag(s12, "POR QUE APRENDER  ·  Razão 2 de 3")

txt(s12, "Big O é a língua franca da engenharia de software", 0.35, 0.45, 9.3, 0.5,
    size=20, bold=True, color=C_WHITE, font=FONT_HEAD)
txt(s12, "Saber dizer a complexidade separa o dev médio do dev sênior — em code review, entrevistas e arquitetura.",
    0.35, 0.95, 9.3, 0.38, size=12.5, italic=True,
    color=RGBColor(0xA0,0xCC,0xE8))

# Bad solution left
add_rect(s12, 0.35, 1.42, 4.6, 0.3, fill=RGBColor(0x6B,0x10,0x10))
txt(s12, "❌  O(n²)  — Matchmaking ingênuo", 0.45, 1.42, 4.4, 0.3,
    size=11.5, bold=True, color=C_RED, font=FONT_HEAD)

code_block(s12, [
    "def encontrar_par_ruim(jogadores):",
    "    for j1 in jogadores:     # O(n)",
    "        for j2 in jogadores: # O(n²) total",
    "            if compativel(j1, j2):",
    "                return j1, j2",
    "# 10.000 jogadores → 100.000.000 ops!",
], 0.35, 1.74, 4.6, 1.78, highlight_lines=[5])

# Good solution right
add_rect(s12, 5.1, 1.42, 4.55, 0.3, fill=RGBColor(0x06,0x3A,0x20))
txt(s12, "✅  O(n log n)  — Matchmaking eficiente", 5.2, 1.42, 4.3, 0.3,
    size=11.5, bold=True, color=C_GREEN, font=FONT_HEAD)

code_block(s12, [
    "# 1. Ordena por ranking — O(n log n)",
    "jogadores_ord = sorted(",
    "    jogadores, key=lambda j: j.ranking)",
    "",
    "# 2. Busca binária pelo mais próximo",
    "pos = bisect.bisect_left(rankings, r)",
    "# O(log n) por jogador → O(n log n) total",
], 5.1, 1.74, 4.55, 1.78)

# Bottom callout
add_rect(s12, 0.35, 3.65, 9.3, 0.85, fill=RGBColor(0x06,0x3A,0x56))
txt_ml(s12, [
    ("Em entrevista técnica:", True, C_GOLD),
    ("  Entrevistador: \"Qual a complexidade?\" — ", False, C_WHITE),
], 0.5, 3.72, 4.5, 0.72, size=13)
txt_ml(s12, [
    ('❌  "Funciona, testei aqui"  →  entrevista encerrada', True, C_RED),
    ('✅  "O(n log n): ordenação + busca binária"  →  próxima fase', True, C_GREEN),
], 4.8, 3.72, 5.1, 0.72, size=12.5)

txt(s12, "10.000 jogadores:  ❌ 100.000.000 ops  vs  ✅ ~130.000 ops",
    0.35, 4.6, 9.3, 0.38, size=13, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 13  —  RAZÃO 5: Transferência de raciocínio
# ════════════════════════════════════════════════════════════════════════
s13 = slide(dark=True)
add_logo(s13)
section_tag(s13, "POR QUE APRENDER  ·  Razão 3 de 3")

txt(s13, "Os algoritmos ensinam a PENSAR, não só a programar", 0.35, 0.45, 9.3, 0.5,
    size=20, bold=True, color=C_WHITE, font=FONT_HEAD)
txt(s13, "O raciocínio da Busca Binária — eliminar metade a cada passo — aparece em todo lugar.",
    0.35, 0.95, 9.3, 0.35, size=12.5, italic=True,
    color=RGBColor(0xA0,0xCC,0xE8))

# Left: git bisect
add_rect(s13, 0.35, 1.38, 4.6, 0.3, fill=C_TEAL2)
txt(s13, "🔍  git bisect — achar o commit com bug", 0.45, 1.38, 4.4, 0.3,
    size=11.5, bold=True, color=C_WHITE, font=FONT_HEAD)

code_block(s13, [
    "# 100 commits → máx. 7 testes!",
    "def git_bisect(commits, tem_bug):",
    "    esq, dir = 0, len(commits) - 1",
    "    while esq < dir:",
    "        meio = (esq + dir) // 2",
    "        if tem_bug(commits[meio]):",
    "            dir = meio   # bug já existe",
    "        else:",
    "            esq = meio + 1  # sem bug",
    "    return commits[esq]  # culpado!",
], 0.35, 1.7, 4.6, 2.38)

# Right: timeout optimization
add_rect(s13, 5.1, 1.38, 4.55, 0.3, fill=C_TEAL2)
txt(s13, "⚙️  Busca binária para otimizar timeout", 5.2, 1.38, 4.3, 0.3,
    size=11.5, bold=True, color=C_WHITE, font=FONT_HEAD)

code_block(s13, [
    "# Menor timeout viável num sistema",
    "def menor_timeout(sistema, mn=100, mx=5000):",
    "    while mn < mx:",
    "        meio = (mn + mx) // 2",
    "        if sistema.ok(meio):",
    "            mx = meio    # tenta menor",
    "        else:",
    "            mn = meio + 1  # aumenta",
    "    return mn  # mínimo que funciona",
], 5.1, 1.7, 4.55, 2.38)

# Bottom insight
add_rect(s13, 0.35, 4.22, 9.3, 0.95, fill=RGBColor(0x06,0x2A,0x1A))
add_rect(s13, 0.35, 4.22, 0.06, 0.95, fill=C_GREEN)
txt_ml(s13, [
    ("Mesmo raciocínio. Contextos diferentes.", True, C_GREEN),
    ("Busca Binária → git bisect → otimização de sistema → A/B testing → debugging.", False, C_WHITE),
    ("Isso é transferência de conhecimento — o objetivo real de qualquer formação técnica.", False, RGBColor(0xA0,0xCC,0xE8)),
], 0.55, 4.3, 9.0, 0.8, size=12.5)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 14  —  RESUMO
# ════════════════════════════════════════════════════════════════════════
s14 = slide(dark=True)
add_logo(s14)
section_tag(s14, "FECHAMENTO  ·  O que aprendemos")

txt(s14, "Resumo da aula", 0.35, 0.45, 9, 0.5,
    size=24, bold=True, color=C_WHITE, font=FONT_HEAD)

summary = [
    ("🔍", "Busca Linear",   "O(n) — simples, sem pré-requisito",            C_TEAL),
    ("⚡", "Busca Binária",  "O(log n) — rápida, lista deve ser ordenada",   C_GREEN),
    ("🫧", "Bubble Sort",    "O(n²) — didático, não usar em produção",        C_MID),
    ("🔀", "Merge Sort",     "O(n log n) — base do sorted() do Python",       C_TEAL2),
    ("📊", "Big O",          "Linguagem universal para comparar algoritmos",   C_GOLD),
    ("💼", "Mercado",        "Timsort, B-Trees, bisect são o padrão real",     C_GREEN),
]
for i, (icon, title, desc, col) in enumerate(summary):
    xi = 0.35 + (i % 2) * 4.85
    yi = 1.22 + (i // 2) * 1.25
    add_rect(s14, xi, yi, 0.55, 0.85, fill=col)
    txt(s14, icon, xi, yi, 0.55, 0.85, size=20, align=PP_ALIGN.CENTER, color=C_WHITE)
    txt(s14, title, xi+0.65, yi+0.04, 3.8, 0.38,
        size=14, bold=True, color=C_WHITE, font=FONT_HEAD)
    txt(s14, desc, xi+0.65, yi+0.44, 3.8, 0.35, size=12.5,
        color=RGBColor(0xA8,0xCC,0xE8))

add_rect(s14, 0.35, 5.0, 9.3, 0.38, fill=C_TEAL)
txt(s14, '"Ordenar uma vez e buscar com eficiência — esse é o padrão de todo sistema escalável."',
    0.5, 5.02, 9.1, 0.34, size=12.5, bold=True, color=C_BG_DARK,
    align=PP_ALIGN.CENTER, italic=True)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 15  —  RECURSOS
# ════════════════════════════════════════════════════════════════════════
s15 = slide(dark=True)
add_logo(s15)
section_tag(s15, "PRÓXIMOS PASSOS  ·  Para continuar aprendendo")

txt(s15, "Quer se aprofundar?", 0.35, 0.45, 9, 0.5,
    size=22, bold=True, color=C_WHITE, font=FONT_HEAD)

resources = [
    ("📺", "Visualgo.net", "Visualizações animadas de algoritmos — gratuito"),
    ("💻", "LeetCode — Easy", "Problemas: Two Sum, Binary Search, Sort Colors"),
    ("📘", "Python Docs", "bisect, heapq, sorted() com key= — documentação oficial"),
    ("🎥", "CS50 Harvard", "Aula de algoritmos no YouTube — grátis e legendado em PT"),
]
for i, (icon, title, desc) in enumerate(resources):
    yi = 1.15 + i * 0.98
    add_rect(s15, 0.4, yi, 9.2, 0.85, fill=RGBColor(0x0F,0x2A,0x40))
    add_rect(s15, 0.4, yi, 0.06, 0.85, fill=C_TEAL)
    txt(s15, icon, 0.5, yi+0.12, 0.6, 0.6, size=24, align=PP_ALIGN.CENTER, color=C_WHITE)
    txt(s15, title, 1.2, yi+0.1, 3.0, 0.38, size=14, bold=True, color=C_TEAL, font=FONT_HEAD)
    txt(s15, desc, 1.2, yi+0.48, 8.0, 0.32, size=12.5, color=RGBColor(0xAA,0xCC,0xEE))

txt(s15, "🏋️  Desafio: implemente a Busca Binária sem consultar o slide. Depois cronometro!",
    0.4, 5.1, 9.2, 0.32, size=13, bold=True, color=C_GOLD, italic=True)

# ════════════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════════════
OUT = "algoritmos_busca_ordenacao_v2.pptx"
prs.save(OUT)
print(f"Saved: {OUT}  ({prs.slides.__len__()} slides)")
